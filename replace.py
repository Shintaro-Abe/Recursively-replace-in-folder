import os
import re
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

def is_text_file(file_path):
    """
    Check if a file is likely to be a text file based on its content.
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            file.read(1024)  # Try to read the first 1024 bytes
        return True
    except UnicodeDecodeError:
        return False

def replace_in_text_file(file_path, old_string, new_string):
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    
    new_content, count = re.subn(re.escape(old_string), new_string, content)
    
    if count > 0:
        with open(file_path, 'w', encoding='utf-8') as file:
            file.write(new_content)
    
    return count

def replace_in_word(file_path, old_string, new_string):
    doc = Document(file_path)
    count = 0
    
    for para in doc.paragraphs:
        if old_string in para.text:
            inline = para.runs
            for i in range(len(inline)):
                if old_string in inline[i].text:
                    text = inline[i].text.replace(old_string, new_string)
                    inline[i].text = text
                    count += 1
    
    if count > 0:
        doc.save(file_path)
    
    return count

def replace_in_excel(file_path, old_string, new_string):
    wb = load_workbook(file_path)
    count = 0
    
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value and old_string in str(cell.value):
                    cell.value = str(cell.value).replace(old_string, new_string)
                    count += 1
    
    if count > 0:
        wb.save(file_path)
    
    return count

def replace_in_powerpoint(file_path, old_string, new_string):
    prs = Presentation(file_path)
    count = 0
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and old_string in shape.text:
                shape.text = shape.text.replace(old_string, new_string)
                count += 1
    
    if count > 0:
        prs.save(file_path)
    
    return count

def process_files(folder_path, old_string, new_string):
    total_files = 0
    total_replacements = 0
    
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            file_path = os.path.join(root, file)
            count = 0
            
            if file.endswith('.docx'):
                count = replace_in_word(file_path, old_string, new_string)
            elif file.endswith('.xlsx'):
                count = replace_in_excel(file_path, old_string, new_string)
            elif file.endswith('.pptx'):
                count = replace_in_powerpoint(file_path, old_string, new_string)
            elif is_text_file(file_path):
                count = replace_in_text_file(file_path, old_string, new_string)
            
            if count > 0:
                total_files += 1
                total_replacements += count
    
    return total_files, total_replacements

def main():
    folder_path = input("フォルダのパス: ")
    old_string = input("置換前: ")
    new_string = input("置換後: ")
    
    total_files, total_replacements = process_files(folder_path, old_string, new_string)
    
    print(f"{total_files}ファイル、{total_replacements}箇所を置換完了")

if __name__ == "__main__":
    main()